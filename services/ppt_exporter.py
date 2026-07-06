import os
from pptx import Presentation
from pptx.util import Inches
from config import Config

try:
    from PIL import Image
except ImportError:
    Image = None


class PowerPointExporter:
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
    PPT_EXTS = {".ppt", ".pptx"}

    # 圖片放大比例
    # 1.0 = 跟 PowerPoint placeholder 一樣大
    # 1.25 = 放大 25%
    # 1.5 = 放大 50%
    IMAGE_SCALE = 1.25

    DEBUG_IMAGE_PLACEHOLDER = True

    SECTION_LABELS = {
        "problem": "Problem",
        "problem_ppt": "Problem PPT",
        "action_taken": "Action Taken",
        "container": "Container",
        "root_cause": "Root Cause",
        "solution": "Solution",
        "implementation": "Implementation",
        "monitoring": "Monitoring"
    }

    def __init__(self, export_folder):
        self.export_folder = export_folder
        os.makedirs(self.export_folder, exist_ok=True)

    def export(self, application, attachments):
        template_path = getattr(Config, "PPT_TEMPLATE_PATH", "")

        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        attachments_by_section = self._group_attachments_by_section(attachments)

        mapping = self._build_placeholder_mapping(
            application=application,
            attachments_by_section=attachments_by_section
        )

        self._replace_text_placeholders_in_presentation(prs, mapping)

        self._replace_image_placeholders(
            prs=prs,
            attachments_by_section=attachments_by_section
        )

        output_path = os.path.join(
            self.export_folder,
            f"{self._value(application, 'request_no')}.pptx"
        )

        prs.save(output_path)

        return output_path

    # --------------------------------------------------
    # Basic helpers
    # --------------------------------------------------

    def _clean_text(self, value):
        if value is None:
            return ""

        text = str(value)

        # Fix PowerPoint _x000D_ issue
        text = text.replace("\r\n", "\n")
        text = text.replace("\r", "\n")
        text = text.replace("_x000D_", "\n")

        return text

    def _value(self, row, key, default=""):
        try:
            value = row[key]
            return value if value is not None else default
        except Exception:
            return default

    def _get_extension(self, file_path):
        if not file_path or "." not in file_path:
            return ""

        return os.path.splitext(file_path)[1].lower()

    def _is_image_attachment(self, attachment):
        file_path = self._value(attachment, "file_path")
        ext = self._get_extension(file_path)
        return ext in self.IMAGE_EXTS

    def _is_ppt_attachment(self, attachment):
        file_path = self._value(attachment, "file_path")
        ext = self._get_extension(file_path)
        return ext in self.PPT_EXTS

    def _get_attachment_full_path(self, attachment):
        return os.path.join(
            Config.UPLOAD_FOLDER,
            self._value(attachment, "file_path")
        )

    def _group_attachments_by_section(self, attachments):
        result = {}

        for section_key in self.SECTION_LABELS.keys():
            result[section_key] = []

        for attachment in attachments:
            section_name = self._value(attachment, "section_name", "others")

            if section_name not in result:
                result[section_name] = []

            result[section_name].append(attachment)

        return result

    # --------------------------------------------------
    # Placeholder mapping
    # --------------------------------------------------

    def _build_placeholder_mapping(self, application, attachments_by_section):
        problem_input_mode = self._value(application, "problem_input_mode", "manual")

        if problem_input_mode == "ppt":
            problem_description = "Problem information was provided by an uploaded PowerPoint file."
            problem_timeline = ""
            problem_attachment_list = self._build_attachment_text_list(
                attachments_by_section.get("problem_ppt", [])
            )
        else:
            problem_description = self._clean_text(
                self._value(application, "problem_description")
            )
            problem_timeline = self._clean_text(
                self._value(application, "problem_timeline")
            )
            problem_attachment_list = self._build_attachment_text_list(
                attachments_by_section.get("problem", [])
            )

        mapping = {
            "{{REQUEST_NO}}": self._clean_text(self._value(application, "request_no")),
            "{{APPLY_DATE}}": self._clean_text(self._value(application, "apply_date")),

            "{{TITLE}}": self._clean_text(self._value(application, "title")),
            "{{IN_DN}}": self._clean_text(self._value(application, "in_dn")),
            "{{CREATE_DATE}}": self._clean_text(self._value(application, "create_date")),
            "{{CLOSE_DATE}}": self._clean_text(self._value(application, "close_date")),
            "{{MACHINE_OR_TOOL}}": self._clean_text(self._value(application, "machine_or_tool")),
            "{{MODULE_NAME}}": self._clean_text(self._value(application, "module_name")),
            "{{DEPARTMENT}}": self._clean_text(self._value(application, "department")),
            "{{AUTHOR}}": self._clean_text(self._value(application, "author")),

            "{{PROBLEM_INPUT_MODE}}": self._clean_text(problem_input_mode),
            "{{PROBLEM_DESCRIPTION}}": problem_description,
            "{{PROBLEM_TIMELINE}}": problem_timeline,
            "{{PROBLEM_ATTACHMENT_LIST}}": problem_attachment_list,

            "{{ACTION_TAKEN}}": self._clean_text(self._value(application, "action_taken")),
            "{{IMPACT}}": self._clean_text(self._value(application, "impact")),
            "{{CONTAINER}}": self._clean_text(self._value(application, "container")),
            "{{NEED_HELP}}": self._clean_text(self._value(application, "need_help")),

            "{{ROOT_CAUSE_DESCRIPTION}}": self._clean_text(
                self._value(application, "root_cause_description")
            ),
            "{{ROOT_CAUSE_POSSIBLE_CAUSE}}": self._clean_text(
                self._value(application, "root_cause_possible_cause")
            ),
            "{{ROOT_CAUSE_TROUBLESHOOTING_TIMELINE}}": self._clean_text(
                self._value(application, "root_cause_troubleshooting_timeline")
            ),

            "{{SOLUTION}}": self._clean_text(self._value(application, "solution")),
            "{{IMPLEMENTATION}}": self._clean_text(self._value(application, "implementation")),
            "{{MONITORING}}": self._clean_text(self._value(application, "monitoring")),
        }

        for section_key in self.SECTION_LABELS.keys():
            placeholder = "{{" + section_key.upper() + "_ATTACHMENT_LIST}}"

            if placeholder == "{{PROBLEM_ATTACHMENT_LIST}}":
                continue

            mapping[placeholder] = self._build_attachment_text_list(
                attachments_by_section.get(section_key, [])
            )

        mapping["{{ALL_ATTACHMENT_LIST}}"] = self._build_all_attachment_text_list(
            attachments_by_section
        )

        self._add_image_attachment_text_mapping(
            mapping=mapping,
            attachments_by_section=attachments_by_section
        )

        return mapping

    def _add_image_attachment_text_mapping(self, mapping, attachments_by_section):
        max_count = getattr(Config, "MAX_IMAGE_COUNT", 10)

        for section_key in self.SECTION_LABELS.keys():
            section_prefix = section_key.upper()

            image_attachments = [
                attachment for attachment in attachments_by_section.get(section_key, [])
                if self._is_image_attachment(attachment)
            ]

            for index in range(1, max_count + 1):
                remark_key = f"{{{{{section_prefix}_IMAGE_{index}_REMARK}}}}"
                filename_key = f"{{{{{section_prefix}_IMAGE_{index}_FILENAME}}}}"
                attachment_no_key = f"{{{{{section_prefix}_IMAGE_{index}_NO}}}}"

                if index <= len(image_attachments):
                    attachment = image_attachments[index - 1]

                    mapping[remark_key] = self._clean_text(
                        self._value(attachment, "remark")
                    )

                    mapping[filename_key] = self._clean_text(
                        self._value(attachment, "original_file_name")
                        or self._value(attachment, "file_path")
                    )

                    mapping[attachment_no_key] = self._clean_text(
                        self._value(attachment, "attachment_no")
                    )
                else:
                    mapping[remark_key] = ""
                    mapping[filename_key] = ""
                    mapping[attachment_no_key] = ""

    def _build_attachment_text_list(self, attachments):
        if not attachments:
            return "No attachment."

        lines = []

        for attachment in attachments:
            attachment_no = self._clean_text(self._value(attachment, "attachment_no"))
            original_file_name = self._clean_text(self._value(attachment, "original_file_name"))
            file_path = self._clean_text(self._value(attachment, "file_path"))
            file_type = self._clean_text(self._value(attachment, "file_type"))
            remark = self._clean_text(self._value(attachment, "remark"))

            lines.append(
                f"Attachment {attachment_no}: {original_file_name or file_path}\n"
                f"Type: {file_type or ''}\n"
                f"Remark: {remark or ''}"
            )

        return "\n\n".join(lines)

    def _build_all_attachment_text_list(self, attachments_by_section):
        lines = []

        for section_key, attachments in attachments_by_section.items():
            if not attachments:
                continue

            section_title = self.SECTION_LABELS.get(section_key, section_key)

            lines.append(f"[{section_title}]")

            for attachment in attachments:
                attachment_no = self._clean_text(self._value(attachment, "attachment_no"))
                original_file_name = self._clean_text(self._value(attachment, "original_file_name"))
                file_path = self._clean_text(self._value(attachment, "file_path"))
                file_type = self._clean_text(self._value(attachment, "file_type"))
                remark = self._clean_text(self._value(attachment, "remark"))

                lines.append(
                    f"Attachment {attachment_no}: {original_file_name or file_path}\n"
                    f"Type: {file_type or ''}\n"
                    f"Remark: {remark or ''}"
                )

            lines.append("")

        if not lines:
            return "No attachment."

        return "\n".join(lines)

    # --------------------------------------------------
    # Replace text placeholders
    # --------------------------------------------------

    def _replace_text_placeholders_in_presentation(self, prs, mapping):
        for slide in prs.slides:
            self._replace_text_placeholders_in_shapes(slide.shapes, mapping)

    def _replace_text_placeholders_in_shapes(self, shapes, mapping):
        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                if shape.has_text_frame:
                    self._replace_text_placeholders_in_text_frame(
                        shape.text_frame,
                        mapping
                    )

            if getattr(shape, "has_table", False):
                if shape.has_table:
                    self._replace_text_placeholders_in_table(shape.table, mapping)

            if hasattr(shape, "shapes"):
                self._replace_text_placeholders_in_shapes(shape.shapes, mapping)

    def _replace_text_placeholders_in_table(self, table, mapping):
        for row in table.rows:
            for cell in row.cells:
                self._replace_text_placeholders_in_text_frame(
                    cell.text_frame,
                    mapping
                )

    def _replace_text_placeholders_in_text_frame(self, text_frame, mapping):
        for paragraph in text_frame.paragraphs:

            if paragraph.runs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text

                    for key, value in mapping.items():
                        safe_value = self._clean_text(value)

                        if key in new_text:
                            new_text = new_text.replace(key, safe_value)

                    if new_text != original_text:
                        run.text = new_text

            paragraph_text = paragraph.text
            new_paragraph_text = paragraph_text

            for key, value in mapping.items():
                safe_value = self._clean_text(value)

                if key in new_paragraph_text:
                    new_paragraph_text = new_paragraph_text.replace(
                        key,
                        safe_value
                    )

            if new_paragraph_text != paragraph_text:
                paragraph.text = new_paragraph_text

    # --------------------------------------------------
    # Replace image placeholders
    # --------------------------------------------------

    def _replace_image_placeholders(self, prs, attachments_by_section):
        max_count = getattr(Config, "MAX_IMAGE_COUNT", 10)

        image_placeholder_mapping = {}

        for section_key in self.SECTION_LABELS.keys():
            section_prefix = section_key.upper()

            for index in range(1, max_count + 1):
                image_placeholder = f"{{{{{section_prefix}_IMAGE_{index}}}}}"
                image_placeholder_mapping[image_placeholder] = (
                    section_key,
                    index - 1
                )

        if self.DEBUG_IMAGE_PLACEHOLDER:
            print("========== IMAGE PLACEHOLDER DEBUG START ==========")
            print("Sections in attachments_by_section:", list(attachments_by_section.keys()))

        for slide_index, slide in enumerate(prs.slides, start=1):
            if self.DEBUG_IMAGE_PLACEHOLDER:
                print(f"Checking slide {slide_index}")

            for shape_index, shape in enumerate(list(slide.shapes), start=1):
                if not getattr(shape, "has_text_frame", False):
                    continue

                if not shape.has_text_frame:
                    continue

                shape_text = self._clean_text(shape.text).strip()

                if not shape_text:
                    continue

                if self.DEBUG_IMAGE_PLACEHOLDER:
                    print(f"Slide {slide_index}, Shape {shape_index}, Text: {shape_text}")

                matched_placeholder = None

                for placeholder in image_placeholder_mapping.keys():
                    if placeholder in shape_text:
                        matched_placeholder = placeholder
                        break

                if not matched_placeholder:
                    continue

                if self.DEBUG_IMAGE_PLACEHOLDER:
                    print("Matched placeholder:", matched_placeholder)

                section_key, image_index = image_placeholder_mapping[matched_placeholder]

                image_attachments = [
                    attachment for attachment in attachments_by_section.get(section_key, [])
                    if self._is_image_attachment(attachment)
                ]

                if self.DEBUG_IMAGE_PLACEHOLDER:
                    print("Section:", section_key)
                    print("Image attachment count:", len(image_attachments))
                    print("Image index:", image_index)

                if image_index >= len(image_attachments):
                    if self.DEBUG_IMAGE_PLACEHOLDER:
                        print("No image attachment for this placeholder.")
                    shape.text = ""
                    continue

                attachment = image_attachments[image_index]
                image_path = self._get_attachment_full_path(attachment)

                if self.DEBUG_IMAGE_PLACEHOLDER:
                    print("Image path:", image_path)
                    print("Image exists:", os.path.exists(image_path))
                    print("Placeholder left:", shape.left)
                    print("Placeholder top:", shape.top)
                    print("Placeholder width:", shape.width)
                    print("Placeholder height:", shape.height)

                if not os.path.exists(image_path):
                    if self.DEBUG_IMAGE_PLACEHOLDER:
                        print("Image file not found. Skip.")
                    shape.text = ""
                    continue

                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                if width <= 0 or height <= 0:
                    if self.DEBUG_IMAGE_PLACEHOLDER:
                        print("Invalid placeholder size. Skip.")
                    shape.text = ""
                    continue

                try:
                    element = shape._element
                    element.getparent().remove(element)
                except Exception:
                    shape.text = ""

                scale = self.IMAGE_SCALE

                new_width = int(width * scale)
                new_height = int(height * scale)

                new_left = left - int((new_width - width) / 2)
                new_top = top - int((new_height - height) / 2)

                self._add_picture_fit(
                    slide=slide,
                    image_path=image_path,
                    left=new_left,
                    top=new_top,
                    max_width=new_width,
                    max_height=new_height
                )

                if self.DEBUG_IMAGE_PLACEHOLDER:
                    print("Image inserted successfully.")

        if self.DEBUG_IMAGE_PLACEHOLDER:
            print("========== IMAGE PLACEHOLDER DEBUG END ==========")

    def _add_picture_fit(self, slide, image_path, left, top, max_width, max_height):
        if Image is None:
            slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=max_width
            )
            return

        with Image.open(image_path) as img:
            img_width, img_height = img.size

        if img_width <= 0 or img_height <= 0:
            return

        max_w = int(max_width)
        max_h = int(max_height)

        if max_w <= 0 or max_h <= 0:
            return

        image_ratio = img_width / img_height
        box_ratio = max_w / max_h

        if image_ratio >= box_ratio:
            display_width = max_w
            display_height = int(max_w / image_ratio)
        else:
            display_height = max_h
            display_width = int(max_h * image_ratio)

        new_left = left + int((max_w - display_width) / 2)
        new_top = top + int((max_h - display_height) / 2)

        slide.shapes.add_picture(
            image_path,
            new_left,
            new_top,
            width=display_width,
            height=display_height
        )