import os
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from config import Config
from models.application import Application
from models.attachment import Attachment
from utils.file_helper import FileHelper

class PowerPointImporter:
    
    def __init__(self, template_path: str = Config.PPT_TEMPLATE_PATH):
        self.template_path = template_path
        self.shape_mapping = self._build_shape_mapping()
        
    def _build_shape_mapping(self) -> dict:
        """
        Scan the original template and build a mapping:
        { "ShapeName": "FIELD_NAME" }
        """
        mapping = {}
        if not os.path.exists(self.template_path):
            print(f"[Warning] PPT Template not found at {self.template_path}. Import mapping will fail.")
            return mapping
            
        prs = Presentation(self.template_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    if "{{TITLE}}" in text:
                        mapping[shape.name] = "title"
                    elif "{{IN_DN}}" in text:
                        mapping[shape.name] = "in_dn"
                    elif "{{CREATE_DATE}}" in text:
                        mapping[shape.name] = "create_date"
                    elif "{{CLOSE_DATE}}" in text:
                        mapping[shape.name] = "close_date"
                    elif "{{MACHINE_OR_TOOL}}" in text:
                        mapping[shape.name] = "machine_or_tool"
                    elif "{{MODULE_NAME}}" in text:
                        mapping[shape.name] = "module_name"
                    elif "{{DEPARTMENT}}" in text:
                        mapping[shape.name] = "department"
                    elif "{{AUTHOR}}" in text:
                        mapping[shape.name] = "author"
                    
                    elif "{{PROBLEM_DESCRIPTION}}" in text:
                        mapping[shape.name] = "problem_description"
                    elif "{{PROBLEM_TIMELINE}}" in text:
                        mapping[shape.name] = "problem_timeline"
                    elif "{{ACTION_TAKEN}}" in text:
                        mapping[shape.name] = "action_taken"
                    elif "{{IMPACT}}" in text:
                        mapping[shape.name] = "impact"
                    elif "{{CONTAINMENT}}" in text:
                        mapping[shape.name] = "container"
                    elif "{{NEED_HELP}}" in text:
                        mapping[shape.name] = "need_help"
                        
                    elif "{{ROOT_CAUSE_DESCRIPTION}}" in text:
                        mapping[shape.name] = "root_cause_description"
                    elif "{{ROOT_CAUSE_POSSIBLE_CAUSE}}" in text:
                        mapping[shape.name] = "root_cause_possible_cause"
                    elif "{{ROOT_CAUSE_TROUBLESHOOTING_TIMELINE}}" in text:
                        mapping[shape.name] = "root_cause_troubleshooting_timeline"
                        
                    elif "{{SOLUTION}}" in text:
                        mapping[shape.name] = "solution"
                    elif "{{IMPLEMENTATION}}" in text:
                        mapping[shape.name] = "implementation"
                    elif "{{MONITORING}}" in text:
                        mapping[shape.name] = "monitoring"
        return mapping
        
    def import_ppt(self, uploaded_file, request_no: str, apply_date: str) -> tuple[Application, list]:
        """
        Extract text from shapes matching the mapping, and extract pictures.
        """
        
        file_stream = io.BytesIO(uploaded_file.read())
        prs = Presentation(file_stream)
        
        extracted_data = {}
        attachments = []
        
        attachment_counter = 1
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Text extraction based on shape name
                if shape.has_text_frame and shape.name in self.shape_mapping:
                    field_name = self.shape_mapping[shape.name]
                    extracted_data[field_name] = shape.text.strip()
                    
                # Image extraction
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_blob = shape.image.blob
                    ext = shape.image.ext
                    filename = f"imported_img_{attachment_counter}.{ext}"
                    
                    # Create a dummy file object that FileHelper.save_attachment can use
                    class DummyFile:
                        def __init__(self, blob, filename):
                            self.blob = blob
                            self.filename = filename
                        def save(self, dst):
                            with open(dst, 'wb') as f:
                                f.write(self.blob)
                                
                    dummy_file = DummyFile(image_blob, filename)
                    
                    saved_filename = FileHelper.save_attachment(
                        file=dummy_file,
                        request_no=request_no,
                        section_name="problem", # Default to problem section
                        attachment_no=attachment_counter
                    )
                    
                    attachments.append(
                        Attachment(
                            request_no=request_no,
                            section_name="problem",
                            attachment_no=attachment_counter,
                            file_path=saved_filename,
                            original_file_name=filename,
                            file_type=ext,
                            remark="Extracted from PPT import"
                        )
                    )
                    attachment_counter += 1
                    
        application = Application(
            request_no=request_no,
            apply_date=apply_date,
            
            title=extracted_data.get("title", "Imported Application"),
            in_dn=extracted_data.get("in_dn", ""),
            create_date=extracted_data.get("create_date", ""),
            close_date=extracted_data.get("close_date", ""),
            machine_or_tool=extracted_data.get("machine_or_tool", ""),
            module_name=extracted_data.get("module_name", ""),
            department=extracted_data.get("department", "Unknown"),
            author=extracted_data.get("author", "Unknown"),
            
            content_input_mode="manual", # Treat extracted data as manual input
            
            problem_description=extracted_data.get("problem_description", ""),
            problem_timeline=extracted_data.get("problem_timeline", ""),
            
            action_taken=extracted_data.get("action_taken", ""),
            
            impact=extracted_data.get("impact", ""),
            
            container=extracted_data.get("container", ""),
            
            need_help=extracted_data.get("need_help", ""),
            
            root_cause_description=extracted_data.get("root_cause_description", ""),
            root_cause_possible_cause=extracted_data.get("root_cause_possible_cause", ""),
            root_cause_troubleshooting_timeline=extracted_data.get("root_cause_troubleshooting_timeline", ""),
            
            solution=extracted_data.get("solution", ""),
            
            implementation=extracted_data.get("implementation", ""),
            
            monitoring=extracted_data.get("monitoring", ""),
            
            labels=[]
        )
        
        return application, attachments

ppt_importer = PowerPointImporter()
